#!/usr/bin/env python3
"""Generate the UW Course Finder showcase slide deck as a .pptx.

Mirrors slides/index.html (problem / solution / impact) with the app's refreshed
look — warm cream canvas, a serif display face, and monospace labels — for
submission to the UW-IT Student Innovation Lab Community Project Showcase.

    pip install python-pptx
    python3 scripts/build_deck.py   # -> slides/UW-Course-Finder.pptx

Fonts: the live app uses Newsreader / Hanken Grotesk / Space Mono (web fonts).
PowerPoint can't embed those reliably, so the deck uses widely available
stand-ins — Georgia (serif), Calibri (sans), Consolas (mono) — that preserve the
same editorial feel on any machine.
"""
from __future__ import annotations

import os

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

# Palette (matches src/index.css)
CREAM = RGBColor(0xF4, 0xF1, 0xEA)
PURPLE = RGBColor(0x4B, 0x2E, 0x83)
PURPLE_D = RGBColor(0x3D, 0x25, 0x69)
GOLD = RGBColor(0xB7, 0xA5, 0x7A)
GOLD_D = RGBColor(0x85, 0x75, 0x4D)
GOLD_BG = RGBColor(0xF3, 0xEE, 0xE1)
INK = RGBColor(0x24, 0x1C, 0x3B)
T2 = RGBColor(0x5B, 0x54, 0x69)
MUTED = RGBColor(0x8C, 0x84, 0x99)
LINE = RGBColor(0xE6, 0xE1, 0xD6)
CANVAS = RGBColor(0xFA, 0xF8, 0xF3)
CARD = RGBColor(0xFF, 0xFF, 0xFF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LAVENDER = RGBColor(0xD7, 0xCF, 0xE6)
PILL_DONE = RGBColor(0xF3, 0xEE, 0xFB)
PILL_PLAN = RGBColor(0xF6, 0xF0, 0xE2)
NODE_PLAN = RGBColor(0xFB, 0xF6, 0xEA)
NODE_AV_LINE = RGBColor(0xDD, 0xD6, 0xC8)
EDGE = RGBColor(0xCF, 0xC8, 0xB9)

SERIF = "Georgia"
SANS = "Calibri"
MONO = "Consolas"

HERE = os.path.dirname(os.path.abspath(__file__))
ROOT = os.path.dirname(HERE)
OUT = os.path.join(ROOT, "slides", "UW-Course-Finder.pptx")

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def slide(bg=CREAM):
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    r.fill.solid()
    r.fill.fore_color.rgb = bg
    r.line.fill.background()
    r.shadow.inherit = False
    return s


def text(s, l, t, w, h, runs, size=18, color=INK, bold=False, align=PP_ALIGN.LEFT,
         anchor=MSO_ANCHOR.TOP, spacing=1.08, font=SANS):
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    if isinstance(runs, str):
        runs = [[(runs, {})]]
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = spacing
        p.space_after = Pt(6)
        for txt, o in para:
            run = p.add_run()
            run.text = txt
            f = run.font
            f.size = Pt(o.get("size", size))
            f.bold = o.get("bold", bold)
            f.italic = o.get("italic", False)
            f.color.rgb = o.get("color", color)
            f.name = o.get("font", font)
    return tb


def card(s, l, t, w, h, fill=CARD, line=LINE, radius=0.06, accent=None):
    shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(1)
    try:
        shp.adjustments[0] = radius
    except Exception:
        pass
    shp.shadow.inherit = False
    if accent:
        bar = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l + 0.26), Inches(t + 0.3), Inches(0.55), Inches(0.1))
        bar.fill.solid()
        bar.fill.fore_color.rgb = accent
        bar.line.fill.background()
        bar.shadow.inherit = False
    return shp


def centered(shape, runs, size, color, font=SANS, spacing=1.0):
    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        p.line_spacing = spacing
        for txt, o in para:
            run = p.add_run()
            run.text = txt
            run.font.size = Pt(o.get("size", size))
            run.font.bold = o.get("bold", False)
            run.font.color.rgb = o.get("color", color)
            run.font.name = o.get("font", font)


def kicker(s, label, color=GOLD_D):
    text(s, 0.92, 0.62, 11.5, 0.5, label.upper(), size=12.5, color=color, bold=True, font=MONO)


def heading(s, title, color=INK):
    text(s, 0.9, 1.0, 11.6, 1.3, title, size=36, color=color, bold=False, spacing=1.0, font=SERIF)


def bullets(s, l, t, w, h, items, size=18, gap=12):
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = 1.14
        p.space_after = Pt(gap)
        lead = p.add_run()
        lead.text = "▸  "
        lead.font.size = Pt(size)
        lead.font.bold = True
        lead.font.color.rgb = GOLD_D
        lead.font.name = SANS
        for txt, o in (item if isinstance(item, list) else [(item, {})]):
            run = p.add_run()
            run.text = txt
            run.font.size = Pt(o.get("size", size))
            run.font.bold = o.get("bold", False)
            run.font.italic = o.get("italic", False)
            run.font.color.rgb = o.get("color", INK)
            run.font.name = SANS


def footer(s, num, light=False):
    c = LAVENDER if light else MUTED
    text(s, 0.92, 6.95, 6, 0.4, "UW Course Finder", size=11, color=c, bold=True)
    text(s, 11.0, 6.95, 1.5, 0.4, f"{num} / 9", size=11, color=c, align=PP_ALIGN.RIGHT, font=MONO)


def connector(s, x1, y1, x2, y2):
    cxn = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    cxn.line.color.rgb = EDGE
    cxn.line.width = Pt(1.5)
    cxn.shadow.inherit = False
    return cxn


def node(s, l, t, code, title, kind):
    fill = PURPLE if kind == "done" else (NODE_PLAN if kind == "plan" else WHITE)
    border = PURPLE if kind == "done" else (GOLD if kind == "plan" else NODE_AV_LINE)
    code_col = WHITE if kind == "done" else PURPLE
    title_col = RGBColor(0xE5, 0xDD, 0xF2) if kind == "done" else T2
    box = card(s, l, t, 1.55, 0.62, fill=fill, line=border, radius=0.16)
    box.line.width = Pt(1.5)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.1)
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.04)
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = code
    r.font.size = Pt(11); r.font.bold = True; r.font.color.rgb = code_col; r.font.name = MONO
    p2 = tf.add_paragraph()
    r2 = p2.add_run(); r2.text = title
    r2.font.size = Pt(8); r2.font.color.rgb = title_col; r2.font.name = SANS


# ---- Slide 1: Title (echoes the app's degree hero) ------------------------
s = slide(PURPLE)
blob = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.6), Inches(-1.4), Inches(4.2), Inches(4.2))
blob.fill.solid(); blob.fill.fore_color.rgb = GOLD
blob.line.fill.background(); blob.shadow.inherit = False
kicker(s, "UW-IT Student Innovation Lab · Community Project Showcase", color=GOLD)
text(s, 0.9, 1.7, 8.8, 1.8, "UW Course Finder", size=62, color=WHITE, bold=False, font=SERIF)
text(s, 0.92, 3.45, 8.6, 0.8, "Visualize a quarter-by-quarter path to your UW degree.",
     size=22, color=LAVENDER)
text(s, 0.92, 4.7, 8.6, 0.6, "Suraj Vaghela  ·  surajv23@uw.edu", size=18, color=WHITE)
badge = card(s, 0.92, 5.45, 4.0, 0.55, fill=PURPLE_D, line=GOLD, radius=0.5)
centered(badge, [[("Built on 14,171 real UW courses", {})]], 13, WHITE, font=MONO)
# progress-ring motif
ring = s.shapes.add_shape(MSO_SHAPE.DONUT, Inches(9.95), Inches(2.4), Inches(2.5), Inches(2.5))
ring.fill.solid(); ring.fill.fore_color.rgb = GOLD
ring.line.fill.background(); ring.shadow.inherit = False
try:
    ring.adjustments[0] = 0.17
except Exception:
    pass
text(s, 9.95, 3.18, 2.5, 1.0, [[("52%", {"size": 34, "color": WHITE, "bold": True, "font": SERIF})],
                               [("COMPLETE", {"size": 10, "color": LAVENDER, "font": MONO})]],
     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, spacing=1.0)
text(s, 0.92, 6.95, 6, 0.4, "University of Washington", size=11, color=LAVENDER, bold=True)
text(s, 11.0, 6.95, 1.5, 0.4, "1 / 9", size=11, color=LAVENDER, align=PP_ALIGN.RIGHT, font=MONO)

# ---- Slide 2: Problem -----------------------------------------------------
s = slide()
kicker(s, "The problem")
heading(s, "UW has the data. Students still plan by hand.")
text(s, 0.92, 2.1, 11.4, 1.1,
     "MyPlan, DARS, and a 14,000-course catalog all exist — but nothing turns "
     "“what are my requirements and prerequisites?” into “what do I take, and when?”",
     size=20, color=T2, spacing=1.25)
bullets(s, 0.92, 3.6, 11.4, 3.0, [
    "Prerequisite chains are a maze — one wrong order can delay graduation a full year.",
    "Degree requirements live in dense department PDFs.",
    "Sequencing 12+ quarters around prereqs and course offerings is manual guesswork.",
    "Transfer and returning students feel it most.",
], size=19)
footer(s, 2)

# ---- Slide 3: Solution ----------------------------------------------------
s = slide()
kicker(s, "The solution")
heading(s, "Your major + your courses → a clear plan.")
text(s, 0.92, 2.05, 11.4, 0.9,
     "Pick your major and the app builds a prerequisite-valid plan instantly — "
     "then you shape it, quarter by quarter.", size=20, color=T2, spacing=1.2)
for i, (t_, d_) in enumerate([
    ("Plan", "An auto-scheduler seeds a quarter-by-quarter plan that respects every prerequisite."),
    ("Edit", "Browse the catalog and add, remove, or move courses between quarters."),
    ("Explore", "An interactive map shows how each course flows into the ones it unlocks."),
]):
    x = 0.92 + i * (3.62 + 0.27)
    card(s, x, 3.35, 3.62, 2.7, accent=GOLD)
    text(s, x + 0.26, 3.85, 3.1, 0.6, t_, size=22, color=PURPLE, bold=False, font=SERIF)
    text(s, x + 0.26, 4.5, 3.1, 1.4, d_, size=15, color=T2, spacing=1.2)
footer(s, 3)

# ---- Slide 4: Two views (refreshed mockups) -------------------------------
s = slide()
kicker(s, "How it works")
heading(s, "Two views, one personalized plan.")

# Left panel — Quarter Plan
card(s, 0.92, 2.2, 5.7, 4.4)
text(s, 1.2, 2.4, 5.2, 0.5, "Quarter Plan", size=21, color=PURPLE, bold=False, font=SERIF)
text(s, 1.2, 2.95, 5.2, 0.7, "An editable schedule to UW’s 180-credit minimum, with an estimated graduation term.",
     size=13, color=MUTED, spacing=1.15)
plan_cols = [
    ("Autumn", RGBColor(0xC9, 0x74, 0x2E), "CSE 143", "Computer Programming II", "5 cr", "Completed", PILL_DONE, PURPLE),
    ("Winter", RGBColor(0x5B, 0x7F, 0xB0), "CSE 311", "Foundations of Computing", "4 cr", "In plan", PILL_PLAN, GOLD_D),
    ("Spring", RGBColor(0x5E, 0x9E, 0x6F), "CSE 332", "Data Structures", "4 cr", "In plan", PILL_PLAN, GOLD_D),
]
cw = 1.72
for i, (term, accent, code, title, cr, pill, pill_bg, pill_col) in enumerate(plan_cols):
    x = 1.2 + i * (cw + 0.06)
    dot = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(3.78), Inches(0.13), Inches(0.13))
    dot.fill.solid(); dot.fill.fore_color.rgb = accent; dot.line.fill.background(); dot.shadow.inherit = False
    text(s, x + 0.2, 3.7, cw, 0.3, term, size=11, color=INK, bold=True)
    mc = card(s, x, 4.05, cw, 1.25, radius=0.1)
    mc.shadow.inherit = False
    text(s, x + 0.13, 4.13, cw - 0.5, 0.3, code, size=10.5, color=PURPLE, bold=True, font=MONO)
    crb = card(s, x + cw - 0.56, 4.16, 0.46, 0.26, fill=GOLD_BG, line=None, radius=0.3)
    centered(crb, [[(cr, {})]], 8, GOLD_D, font=SANS)
    text(s, x + 0.13, 4.45, cw - 0.26, 0.45, title, size=9.5, color=T2, spacing=1.05)
    pl = card(s, x + 0.13, 4.92, 0.95, 0.26, fill=pill_bg, line=None, radius=0.5)
    centered(pl, [[(pill, {})]], 7.5, pill_col, font=SANS)

# Right panel — Prerequisite Map
card(s, 6.85, 2.2, 5.6, 4.4)
text(s, 7.13, 2.4, 5.1, 0.5, "Prerequisite Map", size=21, color=PURPLE, bold=False, font=SERIF)
text(s, 7.13, 2.95, 5.1, 0.7, "A pannable map — drag nodes, zoom, and trace how courses unlock the next.",
     size=13, color=MUTED, spacing=1.15)
cv = card(s, 7.13, 3.65, 5.04, 2.25, fill=CANVAS, line=RGBColor(0xEC, 0xE7, 0xDC), radius=0.05)
cv.shadow.inherit = False
# edges first (behind nodes)
connector(s, 8.86, 4.06, 9.7, 4.62)
connector(s, 8.86, 5.36, 9.7, 4.78)
connector(s, 11.25, 4.7, 11.5, 4.18)
node(s, 7.3, 3.78, "CSE 143", "Computer Programming II", "done")
node(s, 7.3, 5.08, "MATH 126", "Calculus III", "done")
node(s, 9.7, 4.4, "CSE 311", "Foundations of Computing", "plan")
node(s, 10.7, 3.9, "CSE 332", "Data Structures", "av")
# legend
text(s, 7.13, 6.0, 5.2, 0.4, [[
    ("■ ", {"color": PURPLE}), ("Completed   ", {"color": MUTED, "size": 11}),
    ("■ ", {"color": GOLD}), ("Planned   ", {"color": MUTED, "size": 11}),
    ("■ ", {"color": NODE_AV_LINE}), ("Available later", {"color": MUTED, "size": 11}),
]], size=11)
footer(s, 4)

# ---- Slide 5: Under the hood ----------------------------------------------
s = slide()
kicker(s, "Under the hood")
heading(s, "Grounded in real UW data.")
stats = [("14,171", "real UW Seattle courses"), ("4", "majors + specializations"),
         ("180", "credit plans, prereq-valid"), ("58 KB", "tiny, in-browser, private")]
sw = 2.72
for i, (big, lbl) in enumerate(stats):
    x = 0.92 + i * (sw + 0.18)
    card(s, x, 2.15, sw, 1.7)
    text(s, x + 0.22, 2.35, sw - 0.4, 0.8, big, size=32, color=PURPLE, bold=False, font=SERIF)
    text(s, x + 0.22, 3.2, sw - 0.4, 0.55, lbl, size=13, color=MUTED, spacing=1.05)
bullets(s, 0.92, 4.3, 11.5, 2.3, [
    "Course catalog from the official UW Course Descriptions; CSE, Informatics, Math & ACMS requirements encoded.",
    "An auto-scheduler seeds the plan; everything is editable and saved in your browser.",
    "React + TypeScript with a hand-built map canvas — no heavy frameworks; deploys free as a static site.",
], size=18)
footer(s, 5)

# ---- Slide 6: Landscape & Market ------------------------------------------
s = slide()
kicker(s, "Landscape & market")
heading(s, "Different by design — built for every Husky.")
card(s, 0.92, 2.2, 5.7, 4.05)
text(s, 1.2, 2.4, 5.2, 0.5, "What exists today", size=20, color=PURPLE, bold=False, font=SERIF)
_rows = [
    ("MyPlan", "Builds a schedule — no prereq-valid degree path"),
    ("DARS", "A checklist of what's left — no sequencing"),
    ("prereq-map", "Visual only (archived 2022) — no personalization"),
    ("Us", "Auto-planned + editable + interactive prereq map, per major"),
]
_y = 2.98
for _k, _v in _rows:
    _us = _k == "Us"
    text(s, 1.2, _y, 1.35, 0.4, _k, size=12, color=(PURPLE if _us else GOLD_D), bold=True, font=MONO)
    text(s, 2.55, _y, 3.9, 0.7, _v, size=12.5, color=(INK if _us else T2), bold=_us, spacing=1.1)
    _y += 0.78
card(s, 6.85, 2.2, 5.6, 4.05)
text(s, 7.13, 2.4, 5.1, 0.5, "Who it's for", size=20, color=PURPLE, bold=False, font=SERIF)
bullets(s, 7.13, 3.0, 5.05, 3.0, [
    [("~36,000", {"bold": True, "color": PURPLE}),
     (" UW Seattle undergrads (~60k across campuses) — all plan a degree.", {})],
    [("Highest need: ", {}),
     ("freshmen, transfer students, pre-major applicants", {"bold": True}),
     (" to capacity-capped majors, and major-switchers.", {})],
    "Free, no login — zero-friction; expandable to every major and to peer schools.",
], size=14, gap=14)
footer(s, 6)

# ---- Slide 7: Impact ------------------------------------------------------
s = slide()
kicker(s, "The impact")
heading(s, "Clearer planning, on-time graduation.")
bullets(s, 0.92, 2.15, 11.5, 2.6, [
    [("For students: ", {"bold": True, "color": PURPLE}),
     ("better course decisions, fewer wasted credits, less adviser back-and-forth.", {})],
    [("Extensible by design: ", {"bold": True, "color": PURPLE}),
     ("every major is pure data — a new program is one file, no code.", {})],
    [("A genuine gap: ", {"bold": True, "color": PURPLE}),
     ("UW-IT’s own prereq-map (archived 2022) only ", {}), ("visualized", {"italic": True}),
     (" prerequisites; this ", {}), ("plans", {"italic": True}), (" a personalized degree.", {})],
], size=19, gap=16)
co = card(s, 0.92, 4.85, 11.5, 1.4, fill=RGBColor(0xFC, 0xFB, 0xF7), line=GOLD, radius=0.04)
text(s, 1.25, 5.0, 10.9, 1.1,
     "A natural input to MyPlan and advising — turning static requirements into a "
     "living, personal roadmap for every Husky.", size=18, color=INK, spacing=1.25,
     anchor=MSO_ANCHOR.MIDDLE)
footer(s, 7)

# ---- Slide 8: Roadmap -----------------------------------------------------
s = slide()
kicker(s, "What’s next")
heading(s, "Roadmap.")
bullets(s, 0.92, 2.25, 11.5, 3.5, [
    "Expand to more majors and minors across all three campuses.",
    "Model UW’s full “one-of / all-of” prerequisite logic for exact accuracy.",
    "Live course & offerings data via a MyPlan partnership.",
    "Save, compare, and share plans; mobile-first layout.",
], size=20, gap=16)
footer(s, 8)

# ---- Slide 9: Closing -----------------------------------------------------
s = slide(PURPLE)
kicker(s, "Try it", color=GOLD)
text(s, 0.9, 1.6, 11.6, 1.5, "Plan your path.", size=56, color=WHITE, bold=False, font=SERIF)
text(s, 0.92, 3.4, 11.5, 2.2, [
    [("Live demo:  ", {"bold": True, "color": WHITE}),
     ("svog23.github.io/uwcoursefinder", {"color": LAVENDER})],
    [("Code:  ", {"bold": True, "color": WHITE}),
     ("github.com/SVOG23/UWCourseFinder", {"color": LAVENDER})],
    [("Contact:  ", {"bold": True, "color": WHITE}),
     ("Suraj Vaghela · surajv23@uw.edu", {"color": LAVENDER})],
], size=20, spacing=1.6)
text(s, 0.92, 6.1, 11.5, 0.6,
     "An unofficial, student-built planning aid — not affiliated with or endorsed by UW.",
     size=13, color=LAVENDER)
text(s, 0.92, 6.95, 6, 0.4, "University of Washington", size=11, color=LAVENDER, bold=True)
text(s, 11.0, 6.95, 1.5, 0.4, "9 / 9", size=11, color=LAVENDER, align=PP_ALIGN.RIGHT, font=MONO)

os.makedirs(os.path.dirname(OUT), exist_ok=True)
prs.save(OUT)
print(f"Wrote {OUT} ({os.path.getsize(OUT) // 1024} KB, {len(prs.slides._sldIdLst)} slides)")
